from pptx import Presentation
import pytesseract
from PIL import Image

# 이미지에서 텍스트 추출
def extract_text_from_image(image_path):
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img, lang='kor')  # 한국어로 설정
    return text

# 텍스트를 PPT에 삽입
def insert_text_to_ppt(text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 빈 슬라이드 추가
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = text

    # 파일 저장
    prs.save('output.pptx')

# 사용 예
image_path = '~/Desktop/악보1.png'
text = extract_text_from_image(image_path)
insert_text_to_ppt(text)
